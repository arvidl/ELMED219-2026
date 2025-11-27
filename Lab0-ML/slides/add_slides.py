#!/usr/bin/env python3
"""
Script for å legge til slides i ELMED219-2026 Lab0-ML presentasjonen.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Åpne eksisterende presentasjon
pptx_path = "/Users/arvid/GitHub/ELMED219-2026/Lab0-ML/slides/ELMED219-2026_Lab0-ML_01-Enkle_eksempler-slides.pptx"
prs = Presentation(pptx_path)

# Definer farger
UiB_BLUE = RGBColor(0, 58, 112)  # #003A70
ACCENT_ORANGE = RGBColor(230, 126, 34)

def add_title_slide(title, subtitle=""):
    """Legger til en tittelslide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Tittel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = UiB_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, bullets, note=""):
    """Legger til en innholdsslide med punktliste."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Tittel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = UiB_BLUE
    
    # Innhold
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.level = 0
        p.space_before = Pt(12)
    
    # Notat nederst
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)
    
    return slide

def add_quote_slide(title, quote, source=""):
    """Legger til en slide med sitat."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Tittel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = UiB_BLUE
    
    # Sitat
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(28)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    if source:
        p = tf.add_paragraph()
        p.text = f"— {source}"
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_two_column_slide(title, left_title, left_items, right_title, right_items):
    """Legger til en slide med to kolonner."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Hovedtittel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = UiB_BLUE
    
    # Venstre kolonne
    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.space_before = Pt(8)
    
    # Høyre kolonne
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.space_before = Pt(8)
    
    return slide

# ============================================================================
# SLIDE 2: Læringsmål
# ============================================================================
add_content_slide(
    "Hva skal vi lære i dag?",
    [
        "✅ Forstå hva maskinlæring er",
        "✅ Utforske medisinske datasett (Iris, Diabetes)",
        "✅ Trene en klassifikasjonsmodell (Random Forest)",
        "✅ Evaluere modellens ytelse",
        "✅ Forstå Explainable AI (XAI)",
        "✅ Bygge en interaktiv webapp med Gradio"
    ]
)

# ============================================================================
# SLIDE 3: Hva er maskinlæring?
# ============================================================================
add_quote_slide(
    "Hva er maskinlæring?",
    "Maskinlæring er en gren av kunstig intelligens der datamaskiner lærer mønstre fra data, uten å bli eksplisitt programmert.",
    "Arthur Samuel, 1959"
)

# ============================================================================
# SLIDE 4: ML vs tradisjonell programmering
# ============================================================================
add_two_column_slide(
    "ML vs. Tradisjonell programmering",
    "Tradisjonell:",
    [
        "if blodsukker > 126:",
        "    return 'diabetes'",
        "",
        "Programmerer skriver regler",
        "Eksplisitte terskelverdier"
    ],
    "Maskinlæring:",
    [
        "model.fit(X, y)",
        "prediction = model.predict(new_data)",
        "",
        "Modellen lærer regler fra data",
        "Finner terskler automatisk"
    ]
)

# ============================================================================
# SLIDE 5: ML i medisin
# ============================================================================
add_content_slide(
    "Hvorfor maskinlæring i medisin?",
    [
        "🏥 Diagnostikk — Bildediagnostikk, patologi, ECG-analyse",
        "💊 Behandlingsvalg — Prediksjon av behandlingsrespons",
        "⚠️ Risikovurdering — Diabetes, hjertesykdom, sepsis",
        "🔬 Legemiddelutvikling — Molekyldesign, bivirkningsprediksjon",
        "📊 Pasientmonitorering — Tidlig varsling, ICU-prediksjon"
    ]
)

# ============================================================================
# SLIDE 6: Arbeidsflyt i ML
# ============================================================================
add_content_slide(
    "Arbeidsflyt: Fra data til prediksjon",
    [
        "1️⃣ DATA — Samle inn og organisere data",
        "2️⃣ PREPROSESSERING — Rensing, normalisering, splitting",
        "3️⃣ TRENING — Velg algoritme, tren modell",
        "4️⃣ EVALUERING — Test på usett data, mål ytelse",
        "5️⃣ DEPLOYMENT — Sett modellen i produksjon"
    ],
    "I denne labben fokuserer vi på steg 2-4"
)

# ============================================================================
# SLIDE 7: Iris-datasettet
# ============================================================================
add_content_slide(
    "Vårt første datasett: Iris-blomster",
    [
        "🌸 150 blomster, 3 arter (Setosa, Versicolor, Virginica)",
        "📏 4 egenskaper:",
        "     • Sepal length (begerbladlengde)",
        "     • Sepal width (begerbladbredde)",
        "     • Petal length (kronbladlengde)",
        "     • Petal width (kronbladbredde)",
        "🎯 Klassisk benchmark-datasett (Fisher, 1936)"
    ],
    "Enkelt datasett for å forstå ML-konsepter"
)

# ============================================================================
# SLIDE 8: Visualisering
# ============================================================================
add_content_slide(
    "Visualisering: Kan vi skille artene?",
    [
        "📊 Scatter plot viser at:",
        "",
        "     • Setosa er tydelig separert fra de andre",
        "     • Versicolor og Virginica overlapper noe",
        "     • Petal-målinger skiller bedre enn sepal",
        "",
        "💡 God separasjon → Enklere klassifiseringsoppgave"
    ],
    "[Sett inn scatter plot fra notebook her]"
)

# ============================================================================
# SLIDE 9: Random Forest - Intuisjon
# ============================================================================
add_quote_slide(
    "Random Forest: Mange trær gir bedre svar",
    "Som å spørre 100 eksperter og følge flertallet",
    ""
)

# ============================================================================
# SLIDE 10: Random Forest - Detaljer
# ============================================================================
add_content_slide(
    "Hvordan fungerer Random Forest?",
    [
        "🌳 Ensemble av mange beslutningstrær (f.eks. 100)",
        "🎲 Hvert tre trenes på tilfeldig utvalg av data",
        "🔀 Hvert tre bruker tilfeldig utvalg av egenskaper",
        "🗳️ Prediksjon: Flertallsavstemning mellom trærne",
        "",
        "Fordeler: Robust, unngår overfitting, håndterer støy"
    ]
)

# ============================================================================
# SLIDE 11: Train/Test Split
# ============================================================================
add_two_column_slide(
    "Hvorfor dele data i trening og test?",
    "Treningsdata (75%):",
    [
        "Modellen lærer fra disse",
        "Finner mønstre og regler",
        "Kan 'pugge' dataene"
    ],
    "Testdata (25%):",
    [
        "Holdes helt skjult under trening",
        "Evaluerer på USETT data",
        "Estimerer ytelse i virkeligheten"
    ]
)

# ============================================================================
# SLIDE 12: Overfitting
# ============================================================================
add_content_slide(
    "⚠️ Overfitting: Modellen husker, men forstår ikke",
    [
        "Problem: Modellen lærer treningsdata 'utenat'",
        "",
        "Symptom:",
        "     • Treningsdata: 99% accuracy",
        "     • Testdata: 65% accuracy",
        "",
        "Medisinsk konsekvens:",
        "Modellen fungerer på 'gamle' pasienter, men feiler på nye!"
    ],
    "Løsning: Regularisering, kryss-validering, mer data"
)

# ============================================================================
# SLIDE 13: Iris-resultater
# ============================================================================
add_content_slide(
    "Random Forest på Iris: 100% accuracy! 🎉",
    [
        "Konfusjonsmatrise (testdata, 38 blomster):",
        "",
        "                    Predikert",
        "                 Set   Ver   Vir",
        "Faktisk  Setosa   12    0     0",
        "       Versicolor  0    13    0",
        "       Virginica   0    0    13",
        "",
        "Perfekt klassifisering på testdata!"
    ],
    "NB: Dette er et 'lett' datasett"
)

# ============================================================================
# SLIDE 14: Diabetes-datasettet
# ============================================================================
add_content_slide(
    "Et mer realistisk eksempel: Diabetes-prediksjon",
    [
        "👩 768 kvinner fra Pima-indianerstammen (Arizona, USA)",
        "📅 Data fra 1988 (Pima Indians Diabetes Database)",
        "🎯 Mål: Predikere diabetes basert på 8 egenskaper:",
        "",
        "     Graviditeter, Glukose, Blodtrykk, Hudfold,",
        "     Insulin, BMI, Diabetes Pedigree Function, Alder"
    ]
)

# ============================================================================
# SLIDE 15: Utfordringer
# ============================================================================
add_content_slide(
    "Diabetes er vanskeligere enn Iris!",
    [
        "⚖️ Ubalanserte klasser — 65% friske, 35% diabetikere",
        "❓ Manglende verdier — 0 = missing for glucose, BMI, etc.",
        "🔀 Overlappende klasser — Ingen tydelig separasjon",
        "📉 Støy i målingene — Biologisk variasjon",
        "📅 Gammelt datasett — 1988 → 2026"
    ],
    "Representativt for reelle medisinske data!"
)

# ============================================================================
# SLIDE 16: Evaluering
# ============================================================================
add_content_slide(
    "Evaluering: Mer enn bare 'prosent riktig'",
    [
        "Metrikk          Definisjon                        Verdi",
        "─────────────────────────────────────────────────────────",
        "Accuracy      Andel korrekte                       74%",
        "Precision     Av 'diabetes'-pred, % faktisk syke   65%",
        "Recall        Av faktisk syke, % vi finner         58%",
        "F1-score      Balanse mellom prec. og recall       61%"
    ],
    "Hvilken metrikk er viktigst? Det avhenger av konteksten!"
)

# ============================================================================
# SLIDE 17: Konfusjonsmatrise
# ============================================================================
add_two_column_slide(
    "Konfusjonsmatrise: Hvor feiler modellen?",
    "Matrisen:",
    [
        "             Pred: Frisk  Syk",
        "Faktisk Frisk    TN=98   FP=22",
        "        Syk      FN=29   TP=43",
        "",
        "FP = Falske positive",
        "FN = Falske negative"
    ],
    "Medisinsk perspektiv:",
    [
        "FP (22): Unødvendig bekymring",
        "         og testing",
        "",
        "FN (29): FARLIG! Oversett",
        "         sykdom, forsinket",
        "         behandling"
    ]
)

# ============================================================================
# SLIDE 18: Precision vs Recall
# ============================================================================
add_two_column_slide(
    "Trade-off: Precision vs. Recall",
    "Screening (høy recall):",
    [
        "Mål: Fange ALLE syke",
        "Aksepterer falske positive",
        "Eksempel: Årlig helsesjekk",
        "Terskel: Lav (f.eks. 30%)"
    ],
    "Bekreftende test (høy prec.):",
    [
        "Mål: Være SIKKER",
        "Risikerer å overse noen",
        "Eksempel: Før medisinering",
        "Terskel: Høy (f.eks. 70%)"
    ]
)

# ============================================================================
# SLIDE 19: XAI
# ============================================================================
add_quote_slide(
    "Explainable AI (XAI)",
    "En lege kan ikke stole på en 'black box' som sier 'du har diabetes' uten forklaring.",
    ""
)

# ============================================================================
# SLIDE 20: XAI-metoder
# ============================================================================
add_content_slide(
    "XAI-metoder: Åpne den svarte boksen",
    [
        "🏆 Feature Importance — Hvilke egenskaper er viktigst?",
        "📈 Partial Dependence Plots — Hvordan påvirker én egenskap?",
        "🔍 Permutation Importance — Hvor mye taper vi uten en egenskap?",
        "🧩 SHAP — Hvorfor fikk DENNE pasienten dette svaret?",
        "🎯 LIME — Lokal forklaring for enkeltprediksjoner"
    ]
)

# ============================================================================
# SLIDE 21: Feature Importance
# ============================================================================
add_content_slide(
    "Feature Importance: Hvilke egenskaper er viktigst?",
    [
        "glucose        ████████████████████  (0.25)",
        "bmi            ██████████████        (0.18)",
        "age            ██████████            (0.13)",
        "dpf            ████████              (0.10)",
        "pregnancies    ██████                (0.08)",
        "diastolic      ████                  (0.06)",
        "",
        "💡 Glukose er klart viktigst for diabetes-prediksjon!"
    ]
)

# ============================================================================
# SLIDE 22: PDP
# ============================================================================
add_content_slide(
    "Partial Dependence Plot: Glucose → Risiko",
    [
        "PDP viser hvordan predikert sannsynlighet endres",
        "når vi varierer glucose, alt annet holdt konstant:",
        "",
        "     Glucose < 100:  Lav risiko (~20%)",
        "     Glucose 100-125: Moderat risiko (~40%)",
        "     Glucose > 140:  Høy risiko (~70%+)",
        "",
        "💡 Samsvarer med kliniske grenseverdier!"
    ],
    "[Sett inn PDP-figur fra notebook her]"
)

# ============================================================================
# SLIDE 23: Trustworthy AI
# ============================================================================
add_content_slide(
    "Pålitelig AI i medisin: 6 pilarer",
    [
        "1️⃣ NØYAKTIGHET — Fungerer modellen godt nok?",
        "2️⃣ FORKLARBARHET — Kan vi forstå beslutningene?",
        "3️⃣ RETTFERDIGHET — Fungerer den likt for alle grupper?",
        "4️⃣ ROBUSTHET — Tåler den støy og angrep?",
        "5️⃣ PERSONVERN — Er pasientdata beskyttet?",
        "6️⃣ ANSVARLIGHET — Hvem har ansvaret når AI feiler?"
    ],
    "EU AI Act krever dette for høyrisiko-AI"
)

# ============================================================================
# SLIDE 24: Etikk
# ============================================================================
add_content_slide(
    "Etiske utfordringer i medisinsk AI",
    [
        "🌍 POPULASJONSBIAS — Pima-data → norske pasienter?",
        "⏰ HISTORISK BIAS — 1988-data → 2026-pasienter?",
        "👫 KJØNNSBIAS — Kun kvinner i datasettet",
        "📜 INFORMERT SAMTYKKE — Bør pasienten vite om AI?",
        "⚖️ ANSVAR — Hvem har skylden hvis AI tar feil?"
    ],
    "Viktig: Alltid valider på lokal populasjon!"
)

# ============================================================================
# SLIDE 25: Gradio
# ============================================================================
add_content_slide(
    "Fra notebook til webapp med Gradio",
    [
        "Gradio gjør det enkelt å lage interaktive webapper:",
        "",
        "     gr.Interface(",
        "         fn=predict_diabetes,",
        "         inputs=[Slider('Alder'), Slider('BMI'), ...],",
        "         outputs=Textbox('Prediksjon')",
        "     ).launch()",
        "",
        "✅ Ingen frontend-kode nødvendig!"
    ],
    "[Sett inn skjermbilde av Gradio-appen her]"
)

# ============================================================================
# SLIDE 26: Oppsummering
# ============================================================================
add_content_slide(
    "Oppsummering: Hva har vi lært?",
    [
        "✅ ML lærer mønstre fra data automatisk",
        "✅ Random Forest: Robust ensemble av beslutningstrær",
        "✅ Train/test split: Unngå overfitting, evaluer på usett data",
        "✅ Evaluering: Accuracy, precision, recall, F1, konfusjonsmatrise",
        "✅ XAI: Feature importance, PDP forklarer modellen",
        "✅ Trustworthy AI: 6 pilarer for pålitelig medisinsk AI"
    ]
)

# ============================================================================
# SLIDE 27: Neste steg
# ============================================================================
add_content_slide(
    "Videre læring",
    [
        "📘 Ressurser:",
        "     • Molnar: Interpretable ML (gratis online)",
        "     • James et al.: ISLP (gratis online)",
        "     • scikit-learn.org, gradio.app",
        "",
        "🔜 Neste labs:",
        "     • Lab 1: Nettverksvitenskap og PSN",
        "     • Lab 2: Dyp læring",
        "     • Lab 3: Generativ AI og LLM"
    ],
    "Lykke til med videre læring! 🎓"
)

# ============================================================================
# Lagre presentasjonen
# ============================================================================
prs.save(pptx_path)
print(f"✅ Lagt til {len(prs.slides) - 1} nye slides!")
print(f"   Totalt: {len(prs.slides)} slides")
print(f"   Lagret til: {pptx_path}")

